#!/usr/bin/env python3
"""
Tangible — Generateur de pitch .pptx (version stylisée + effets)
12 slides, charte Tangible, 16:9, notes du présentateur.

Effets réels PowerPoint via injection XML :
  - Drop shadows avec blur (a:outerShdw)
  - Gradients linéaires et radiaux (a:gradFill)
  - Glow orange sur éléments clés (a:glow)
  - Soft edges sur watermarks (a:softEdge)
  - Image alpha transparency (a:alphaModFix)

Usage:
    ./pitch/.venv/bin/python3 pitch/generate_pitch.py
Output:
    pitch/Tangible-Pitch.pptx
"""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree

# ───────────────────────── Paths ────────────────────────
ROOT = Path(__file__).resolve().parent.parent
PITCH = ROOT / "pitch"
LOGOS_SRC = ROOT / "content" / "assets" / "logos"
LOGOS_OUT = PITCH / "logos-cropped"
OUTPUT = PITCH / "Tangible-Pitch.pptx"

# ─────────────────── Charte Tangible ────────────────────
INK          = RGBColor(0x2B, 0x2F, 0x33)
INK_DEEP     = RGBColor(0x15, 0x17, 0x1B)
INK_MID      = RGBColor(0x22, 0x25, 0x29)
INK_SOFT     = RGBColor(0x3A, 0x3E, 0x45)
INK_LINE     = RGBColor(0x4D, 0x52, 0x5A)
FLAME        = RGBColor(0xE6, 0x73, 0x3C)
FLAME_BRIGHT = RGBColor(0xFF, 0x8A, 0x4F)
FLAME_DEEP   = RGBColor(0xC5, 0x58, 0x2A)
FLAME_DARK   = RGBColor(0x8A, 0x3E, 0x1F)
FLAME_FADE   = RGBColor(0x5A, 0x3A, 0x2E)
IVORY        = RGBColor(0xFB, 0xF8, 0xF1)
IVORY_DEEP   = RGBColor(0xF3, 0xEC, 0xDE)
PAPER        = RGBColor(0xFF, 0xFF, 0xFF)
MUTED        = RGBColor(0x9C, 0x97, 0x8D)
MUTED_DARK   = RGBColor(0x5E, 0x5B, 0x53)
MUTED_LIGHT  = RGBColor(0xD0, 0xCC, 0xC0)
SUCCESS      = RGBColor(0x3A, 0xA5, 0x74)
DANGER       = RGBColor(0xD1, 0x4A, 0x3D)

FONT_TITLE = "Georgia"
FONT_BODY  = "Calibri"
FONT_MONO  = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 12

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

# ═══════════════════════════════════════════════════════
#   EFFECTS  (XML injection dans spPr)
# ═══════════════════════════════════════════════════════
def _hex(color: RGBColor) -> str:
    return "{:02X}{:02X}{:02X}".format(color[0], color[1], color[2])

def _get_spPr(shape):
    return shape._element.spPr if hasattr(shape._element, 'spPr') else shape._element.find(qn("p:spPr"))

def _remove_existing_fill(spPr):
    """Remove existing solidFill/gradFill/noFill/blipFill/pattFill."""
    for tag in ("a:solidFill", "a:gradFill", "a:noFill", "a:blipFill", "a:pattFill"):
        for el in spPr.findall(qn(tag)):
            spPr.remove(el)

def _ensure_effect_lst(spPr):
    """Get or create <a:effectLst> inside spPr. Place after fill, before xfrm etc."""
    existing = spPr.find(qn("a:effectLst"))
    if existing is not None:
        # Wipe children to rewrite
        for child in list(existing):
            existing.remove(child)
        return existing
    eff = etree.SubElement(spPr, qn("a:effectLst"))
    return eff

def set_gradient_linear(shape, color1: RGBColor, color2: RGBColor, angle_deg: float = 90):
    """Linear gradient fill. angle_deg: 0=right, 90=down, 180=left, 270=up."""
    spPr = _get_spPr(shape)
    _remove_existing_fill(spPr)
    # Insert after nvSpPr sibling - just append, order rarely matters
    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, col in [(0, color1), (100000, color2)]:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", _hex(col))
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", str(int(angle_deg * 60000)))
    lin.set("scaled", "0")
    tileRect = etree.SubElement(gradFill, qn("a:tileRect"))

def set_gradient_linear_with_alpha(shape, color1: RGBColor, alpha1: int,
                                   color2: RGBColor, alpha2: int, angle_deg: float = 90):
    """Linear gradient with alpha control (0-100000)."""
    spPr = _get_spPr(shape)
    _remove_existing_fill(spPr)
    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, col, a in [(0, color1, alpha1), (100000, color2, alpha2)]:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", _hex(col))
        alpha_el = etree.SubElement(srgb, qn("a:alpha"))
        alpha_el.set("val", str(a))
    lin = etree.SubElement(gradFill, qn("a:lin"))
    lin.set("ang", str(int(angle_deg * 60000)))
    lin.set("scaled", "0")
    etree.SubElement(gradFill, qn("a:tileRect"))

def set_gradient_radial(shape, center_color: RGBColor, edge_color: RGBColor,
                        center_alpha: int = 100000, edge_alpha: int = 0,
                        focal_x: int = 50000, focal_y: int = 50000):
    """Radial gradient. Alpha 0-100000."""
    spPr = _get_spPr(shape)
    _remove_existing_fill(spPr)
    gradFill = etree.SubElement(spPr, qn("a:gradFill"))
    gradFill.set("flip", "none")
    gradFill.set("rotWithShape", "1")
    gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
    for pos, col, a in [(0, center_color, center_alpha),
                         (100000, edge_color, edge_alpha)]:
        gs = etree.SubElement(gsLst, qn("a:gs"))
        gs.set("pos", str(pos))
        srgb = etree.SubElement(gs, qn("a:srgbClr"))
        srgb.set("val", _hex(col))
        alpha_el = etree.SubElement(srgb, qn("a:alpha"))
        alpha_el.set("val", str(a))
    path = etree.SubElement(gradFill, qn("a:path"))
    path.set("path", "circle")
    fillToRect = etree.SubElement(path, qn("a:fillToRect"))
    fillToRect.set("l", str(focal_x))
    fillToRect.set("t", str(focal_y))
    fillToRect.set("r", str(100000 - focal_x))
    fillToRect.set("b", str(100000 - focal_y))
    etree.SubElement(gradFill, qn("a:tileRect"))

def set_shadow(shape, blur_pt: float = 20, dist_pt: float = 6,
               direction_deg: float = 90, color: RGBColor = RGBColor(0,0,0),
               alpha: int = 35000):
    """Drop shadow with soft blur."""
    spPr = _get_spPr(shape)
    eff = _ensure_effect_lst(spPr)
    shdw = etree.SubElement(eff, qn("a:outerShdw"))
    shdw.set("blurRad", str(int(Pt(blur_pt))))
    shdw.set("dist", str(int(Pt(dist_pt))))
    shdw.set("dir", str(int(direction_deg * 60000)))
    shdw.set("algn", "ctr")
    shdw.set("rotWithShape", "0")
    srgb = etree.SubElement(shdw, qn("a:srgbClr"))
    srgb.set("val", _hex(color))
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(alpha))

def set_glow(shape, radius_pt: float = 18, color: RGBColor = FLAME, alpha: int = 55000):
    """Glow effect (halo)."""
    spPr = _get_spPr(shape)
    eff = _ensure_effect_lst(spPr)
    glow = etree.SubElement(eff, qn("a:glow"))
    glow.set("rad", str(int(Pt(radius_pt))))
    srgb = etree.SubElement(glow, qn("a:srgbClr"))
    srgb.set("val", _hex(color))
    alpha_el = etree.SubElement(srgb, qn("a:alpha"))
    alpha_el.set("val", str(alpha))

def set_soft_edge(shape, radius_pt: float = 20):
    """Soft edge (blur border fade)."""
    spPr = _get_spPr(shape)
    eff = _ensure_effect_lst(spPr)
    se = etree.SubElement(eff, qn("a:softEdge"))
    se.set("rad", str(int(Pt(radius_pt))))

def set_picture_alpha(pic, alpha_amt: int = 30000):
    """Apply transparency to a picture (0-100000 where 100000=opaque)."""
    blipFill = pic._element.find(".//" + qn("p:blipFill"))
    if blipFill is None:
        return
    blip = blipFill.find(qn("a:blip"))
    if blip is None:
        return
    # Remove any existing alphaModFix
    for el in blip.findall(qn("a:alphaModFix")):
        blip.remove(el)
    alphaModFix = etree.SubElement(blip, qn("a:alphaModFix"))
    alphaModFix.set("amt", str(alpha_amt))

# ────────────────────── Helpers simples ─────────────────
def autocrop_png(src: Path, dst: Path, pad: int = 8):
    dst.parent.mkdir(parents=True, exist_ok=True)
    img = Image.open(src).convert("RGBA")
    bbox = img.getbbox()
    if not bbox:
        img.save(dst)
        return
    w, h = img.size
    l, t, r, b = bbox
    img.crop((max(0, l - pad), max(0, t - pad),
              min(w, r + pad), min(h, b + pad))).save(dst)

def prepare_logos():
    LOGOS_OUT.mkdir(parents=True, exist_ok=True)
    for png in LOGOS_SRC.glob("*.png"):
        dst = LOGOS_OUT / png.name
        if not dst.exists() or dst.stat().st_mtime < png.stat().st_mtime:
            autocrop_png(png, dst)
    return {p.stem: p for p in LOGOS_OUT.glob("*.png")}

def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _no_border(shp):
    shp.line.fill.background()

def add_rect(slide, x, y, w, h, fill_color=None,
             line_color=None, line_width=None,
             shape=MSO_SHAPE.RECTANGLE,
             radius_pct=None):
    shp = slide.shapes.add_shape(shape, x, y, w, h)
    _no_border(shp)
    if fill_color is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
    if line_color is not None:
        shp.line.color.rgb = line_color
        if line_width is not None:
            shp.line.width = line_width
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE and radius_pct is not None:
        adj = shp.adjustments
        adj[0] = min(0.5, max(0.0, radius_pct))
    return shp

def rounded(slide, x, y, w, h, fill_color=None, line_color=None,
            line_width=None, radius_in=0.12):
    shorter = min(w, h)
    r_pct = 0.08
    if shorter > 0:
        target = Emu(int(Inches(radius_in)))
        r_pct = min(0.5, max(0.02, target / shorter))
    return add_rect(slide, x, y, w, h, fill_color, line_color, line_width,
                    shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius_pct=r_pct)

def add_text(slide, x, y, w, h, text, *,
             font=FONT_BODY, size=18, bold=False, italic=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, space_after=0):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box

def thin_line(slide, x, y, length, vertical=False, thickness_pt=1, color=FLAME):
    thickness = Emu(int(Pt(thickness_pt)))
    if vertical:
        add_rect(slide, x, y, thickness, length, fill_color=color)
    else:
        add_rect(slide, x, y, length, thickness, fill_color=color)

def corner_mark(slide, x, y, size_in=0.22, thickness_pt=1.5, color=FLAME):
    t = Emu(int(Pt(thickness_pt)))
    length = Inches(size_in)
    add_rect(slide, x - length, y, length, t, fill_color=color)
    add_rect(slide, x - t, y, t, length, fill_color=color)

def eyebrow(slide, x, y, text, color=FLAME, with_line=True):
    if with_line:
        thin_line(slide, x, y + Inches(0.13), Inches(0.28), color=color, thickness_pt=1.2)
        add_text(slide, x + Inches(0.42), y, Inches(7), Inches(0.4), text,
                 font=FONT_MONO, size=11, bold=True, color=color,
                 align=PP_ALIGN.LEFT)
    else:
        add_text(slide, x, y, Inches(7), Inches(0.4), text,
                 font=FONT_MONO, size=11, bold=True, color=color)

def big_title(slide, y, text, *, color=INK, size=52, italic=False, align=PP_ALIGN.LEFT):
    add_text(slide, Inches(0.7), y, Inches(12), Inches(1.5),
             text, font=FONT_TITLE, size=size, bold=True, italic=italic,
             color=color, align=align, line_spacing=1.05)

def underline_accent(slide, x, y, length_in=0.8, color=FLAME):
    thin_line(slide, x, y, Inches(length_in), thickness_pt=3, color=color)

def progress_footer(slide, page, total=TOTAL_SLIDES, dark=False):
    fg = MUTED if not dark else RGBColor(0x8A, 0x86, 0x7D)
    add_text(slide, Inches(0.7), Inches(7.08), Inches(3), Inches(0.3),
             "TANGIBLE  ·  SCRUM'INNOV 2026",
             font=FONT_MONO, size=9, bold=True, color=fg)
    dots_start_x = Inches(5.2)
    dot_r = Inches(0.07)
    gap = Inches(0.18)
    for i in range(total):
        dx = dots_start_x + gap * i
        dy = Inches(7.15)
        active = (i + 1) == page
        color = FLAME if active else (RGBColor(0x44, 0x48, 0x4D) if dark else RGBColor(0xCF, 0xC9, 0xB8))
        size = Inches(0.12) if active else dot_r
        offset = (Inches(0.12) - size) / 2 if not active else Emu(0)
        dot = add_rect(slide, dx, dy + offset, size, size,
                       fill_color=color, shape=MSO_SHAPE.OVAL)
        if active:
            set_glow(dot, radius_pt=8, color=FLAME, alpha=70000)
    add_text(slide, Inches(11.5), Inches(7.08), Inches(1.6), Inches(0.3),
             f"{page:02d} / {total:02d}",
             font=FONT_MONO, size=10, bold=True, color=FLAME,
             align=PP_ALIGN.RIGHT)

def watermark_logo(slide, logo_path, size_in=5.0, x_emu=None, y_emu=None,
                   alpha_amt=8000):
    if logo_path is None:
        return
    img = Image.open(logo_path)
    ratio = img.width / img.height
    h = Inches(size_in)
    w = Emu(int(h * ratio))
    if x_emu is None:
        x_emu = SLIDE_W - w + Inches(1.0)
    if y_emu is None:
        y_emu = SLIDE_H - h + Inches(0.5)
    pic = slide.shapes.add_picture(str(logo_path), x_emu, y_emu, width=w, height=h)
    set_picture_alpha(pic, alpha_amt)
    return pic

def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text.strip()

def add_logo_picture(slide, logo_path, x, y, height_in):
    img = Image.open(logo_path)
    ratio = img.width / img.height
    h = Inches(height_in)
    w = Emu(int(h * ratio))
    pic = slide.shapes.add_picture(str(logo_path), x, y, width=w, height=h)
    return pic, w, h

def ambient_glow(slide, cx_emu, cy_emu, radius_in=3.5,
                 color=FLAME, alpha=45000):
    """Radial gradient halo (ellipse semi-transparent)."""
    size = Inches(radius_in * 2)
    x = cx_emu - size // 2
    y = cy_emu - size // 2
    el = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    _no_border(el)
    set_gradient_radial(el, color, color, center_alpha=alpha, edge_alpha=0)
    return el

def full_bg_gradient(slide, color_top: RGBColor, color_bottom: RGBColor,
                     angle_deg=90):
    """Cover entire slide with a gradient rect (goes BEHIND other content)."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _no_border(shp)
    set_gradient_linear(shp, color_top, color_bottom, angle_deg)
    return shp

def full_bg_radial(slide, color_center: RGBColor, color_edge: RGBColor,
                    focal_x=50000, focal_y=35000):
    """Cover slide with radial gradient."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _no_border(shp)
    set_gradient_radial(shp, color_center, color_edge,
                         focal_x=focal_x, focal_y=focal_y)
    return shp

# ═══════════════════════════════════════════════════════
#   SLIDE BUILDERS
# ═══════════════════════════════════════════════════════

def slide_1_cover(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INK_DEEP)

    # Full background gradient (ink_deep → black with orange tint on left)
    full_bg_gradient(s, INK_MID, INK_DEEP, angle_deg=135)
    # Orange ambient glow left
    ambient_glow(s, Inches(2.8), Inches(3.8), radius_in=5.5,
                 color=FLAME, alpha=38000)
    # Secondary softer glow right-bottom
    ambient_glow(s, Inches(11), Inches(6.3), radius_in=3.5,
                 color=FLAME_DEEP, alpha=22000)

    # Corner marks
    corner_mark(s, Inches(13.0), Inches(0.35), size_in=0.3)
    t = Emu(int(Pt(1.5)))
    length = Inches(0.3)
    add_rect(s, Inches(0.35), Inches(7.0 - 0.15), length, t, fill_color=FLAME)
    add_rect(s, Inches(0.35), Inches(7.0 - 0.15 - 0.27), t, Inches(0.3), fill_color=FLAME)

    add_text(s, Inches(0.7), Inches(0.4), Inches(6), Inches(0.35),
             "SCRUM'INNOV 2026  —  PITCH 7 MIN",
             font=FONT_MONO, size=10, bold=True, color=FLAME)
    add_text(s, Inches(7), Inches(0.4), Inches(5.6), Inches(0.35),
             "BUT3 INFORMATIQUE · IUT MONTPELLIER",
             font=FONT_MONO, size=10, bold=True, color=MUTED,
             align=PP_ALIGN.RIGHT)

    # Logo with glow halo
    logo = logos.get("tangible-logo-vertical")
    if logo:
        # Halo behind
        ambient_glow(s, Inches(13.333/2), Inches(2.85), radius_in=2.8,
                     color=FLAME, alpha=40000)
        pic, lw, lh = add_logo_picture(s, logo, Inches(0), Inches(1.35), 3.1)
        s.shapes[-1].left = Inches(13.333/2) - lw // 2

    # Separator orange line
    thin_line(s, Inches(13.333/2 - 0.9), Inches(4.7), Inches(1.8),
              thickness_pt=1.5, color=FLAME)

    # Slogan FR
    slog_fr = add_text(s, Inches(0.7), Inches(4.95), Inches(12), Inches(0.8),
                       "Ne louez plus votre passion. Possédez-la.",
                       font=FONT_TITLE, size=38, bold=True, italic=True,
                       color=IVORY, align=PP_ALIGN.CENTER, line_spacing=1.1)

    add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
             "·  DON'T RENT YOUR PASSION. OWN IT.  ·",
             font=FONT_MONO, size=13, bold=True, color=FLAME,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(7.1), Inches(12), Inches(0.3),
             "15–17 AVRIL 2026",
             font=FONT_MONO, size=10, color=MUTED,
             align=PP_ALIGN.CENTER)

    set_notes(s, """
[0:00-1:00 — HOOK]

"Samedi dernier, j'ai voulu revoir The Matrix. Je l'avais 'acheté' sur iTunes en 2019. 15 €. Verdict : 'indisponible dans votre région'. Mon film a disparu.

Parlons cash : quand on achète un film sur Netflix, Amazon, iTunes, on ne possède rien. On loue. Même quand on paie 15 € pour l'avoir 'pour toujours'.

Résultat : Microsoft a arrêté de vendre des films en 2021 — leurs clients ont perdu leurs bibliothèques. Disney+ retire régulièrement des originaux."

Laisser un beat après "Mon film a disparu."
""")

def slide_2_probleme(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INK_DEEP)

    # Radial gradient ambient top-right
    full_bg_radial(s, INK_MID, INK_DEEP, focal_x=75000, focal_y=30000)

    # Giant "≠" decoration
    add_text(s, Inches(7), Inches(0.5), Inches(6), Inches(6),
             "≠", font=FONT_TITLE, size=380, bold=True, italic=True,
             color=RGBColor(0x26, 0x2A, 0x2F),
             align=PP_ALIGN.CENTER)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "LE PROBLÈME")
    big_title(s, Inches(0.95), "Acheter ≠ Posséder.",
              color=IVORY, size=58, italic=True)
    underline_accent(s, Inches(0.7), Inches(2.05), length_in=1.2)

    add_text(s, Inches(0.7), Inches(2.3), Inches(12), Inches(0.5),
             "4 cas documentés. 15 ans de propriété fictive.",
             font=FONT_BODY, size=15, italic=True, color=MUTED_LIGHT)

    cards = [
        ("2021", "Microsoft Movies",  "fermé",     "Arrêt des ventes. Bibliothèques gelées. Aucun export."),
        ("2023", "Disney+",            "retire",    "~50 titres supprimés. Willow, Crater... disparus."),
        ("2024", "iTunes",             "bloque",    "Changement de région ? Bibliothèque inaccessible."),
        ("∞",    "Netflix & co",       "location",  "3 abonnements / foyer. 300 €/an. Zéro fichier possédé."),
    ]
    card_w = Inches(2.85)
    card_h = Inches(3.4)
    gap = Inches(0.2)
    start_x = Inches(0.7)
    y = Inches(3.25)
    for i, (year, title, suffix, body) in enumerate(cards):
        x = start_x + (card_w + gap) * i
        # Card avec gradient vertical (subtle)
        card = rounded(s, x, y, card_w, card_h, fill_color=INK_MID, radius_in=0.14)
        set_gradient_linear(card,
                            RGBColor(0x2E, 0x32, 0x37),
                            RGBColor(0x1E, 0x21, 0x25),
                            angle_deg=90)
        set_shadow(card, blur_pt=30, dist_pt=8,
                   color=RGBColor(0, 0, 0), alpha=45000)
        # Top orange stripe
        stripe = rounded(s, x, y, card_w, Inches(0.10),
                fill_color=FLAME, radius_in=0.02)
        set_glow(stripe, radius_pt=10, color=FLAME, alpha=55000)
        # Year
        year_tb = add_text(s, x + Inches(0.25), y + Inches(0.35),
                 card_w - Inches(0.5), Inches(0.6),
                 year, font=FONT_MONO, size=16, bold=True, color=FLAME)
        thin_line(s, x + Inches(0.25), y + Inches(0.95),
                  Inches(0.5), thickness_pt=1, color=FLAME)
        add_text(s, x + Inches(0.25), y + Inches(1.15),
                 card_w - Inches(0.5), Inches(0.5),
                 title, font=FONT_TITLE, size=20, bold=True, color=IVORY,
                 line_spacing=1.05)
        add_text(s, x + Inches(0.25), y + Inches(1.6),
                 card_w - Inches(0.5), Inches(0.4),
                 suffix, font=FONT_TITLE, size=15, italic=True, color=FLAME,
                 line_spacing=1.05)
        add_text(s, x + Inches(0.25), y + Inches(2.15),
                 card_w - Inches(0.5), card_h - Inches(2.4),
                 body, font=FONT_BODY, size=12, color=MUTED_LIGHT,
                 line_spacing=1.4)

    progress_footer(s, 2, dark=True)
    set_notes(s, """
[1:00 — PROBLÈME]

Pointer la ligne du temps : Microsoft 2021 → Disney+ 2023 → iTunes 2024 → Netflix toujours.
"Microsoft a arrêté de vendre. Disney+ retire. iTunes bloque. Netflix ne vend même pas. Tous ont une faille."
Marquer une pause sur "≠" du titre pour que le jury le lise.
""")

def slide_3_solution(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)

    # Subtle radial gradient bg
    full_bg_radial(s, PAPER, IVORY_DEEP, focal_x=30000, focal_y=30000)

    # Watermark icon
    icon = logos.get("tangible-icon")
    if icon:
        wm = watermark_logo(s, icon, size_in=6.5,
                            x_emu=Inches(9.6), y_emu=Inches(1.2),
                            alpha_amt=6000)
        if wm: set_soft_edge(wm, radius_pt=12)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "LA SOLUTION")

    logo = logos.get("tangible-logo-horizontal")
    if logo:
        add_logo_picture(s, logo, Inches(0.7), Inches(0.95), 0.85)

    big_title(s, Inches(2.0), "Deux piliers. Une promesse.", size=38)
    underline_accent(s, Inches(0.7), Inches(2.85), length_in=1.5)

    card_w = Inches(5.95)
    card_h = Inches(3.9)
    gap = Inches(0.2)
    y = Inches(3.15)

    # --- PLAYER ---
    c1 = rounded(s, Inches(0.7), y, card_w, card_h, fill_color=PAPER, radius_in=0.2)
    set_gradient_linear(c1, PAPER, RGBColor(0xFA, 0xF6, 0xEA), angle_deg=120)
    set_shadow(c1, blur_pt=32, dist_pt=10, alpha=18000)
    # Badge
    badge1 = rounded(s, Inches(0.95), y + Inches(0.3), Inches(0.8), Inches(0.5),
                     fill_color=IVORY_DEEP, radius_in=0.08)
    add_text(s, Inches(0.95), y + Inches(0.32), Inches(0.8), Inches(0.5),
             "01", font=FONT_MONO, size=14, bold=True, color=FLAME,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.95), y + Inches(0.95), card_w - Inches(0.5), Inches(0.7),
             "Tangible Player", font=FONT_TITLE, size=28, bold=True, color=INK)
    add_text(s, Inches(0.95), y + Inches(1.65), card_w - Inches(0.5), Inches(0.4),
             "Media center local premium",
             font=FONT_BODY, size=14, italic=True, color=MUTED_DARK)
    thin_line(s, Inches(0.95), y + Inches(2.1), Inches(0.7),
              color=FLAME, thickness_pt=1.5)
    player_items = [
        "Bibliothèque chiffrée AES-256",
        "4K HDR · Dolby Vision · sous-titres",
        "Chromecast / AirPlay / DLNA",
        "Profils enfants · biométrie",
        "100 % hors-ligne",
    ]
    for i, item in enumerate(player_items):
        by = y + Inches(2.35 + i * 0.3)
        add_text(s, Inches(1.0), by, card_w - Inches(0.5), Inches(0.3),
                 "✓  " + item, font=FONT_BODY, size=13, color=INK_SOFT)

    # --- STORE accented ---
    x2 = Inches(0.7) + card_w + gap
    c2 = rounded(s, x2, y, card_w, card_h, fill_color=PAPER, radius_in=0.2)
    set_gradient_linear(c2, PAPER, RGBColor(0xFD, 0xF4, 0xEC), angle_deg=120)
    set_shadow(c2, blur_pt=40, dist_pt=12,
               color=FLAME_DEEP, alpha=22000)
    # Top orange stripe with gradient
    stripe2 = rounded(s, x2, y, card_w, Inches(0.14),
            fill_color=FLAME, radius_in=0.04)
    set_gradient_linear(stripe2, FLAME_BRIGHT, FLAME_DEEP, angle_deg=0)
    set_glow(stripe2, radius_pt=14, color=FLAME, alpha=55000)
    # Badge
    badge2 = rounded(s, x2 + Inches(0.25), y + Inches(0.4),
            Inches(0.8), Inches(0.5),
            fill_color=FLAME, radius_in=0.08)
    set_gradient_linear(badge2, FLAME_BRIGHT, FLAME_DEEP, angle_deg=120)
    set_glow(badge2, radius_pt=10, color=FLAME, alpha=45000)
    add_text(s, x2 + Inches(0.25), y + Inches(0.42), Inches(0.8), Inches(0.5),
             "02", font=FONT_MONO, size=14, bold=True, color=IVORY,
             align=PP_ALIGN.CENTER)
    add_text(s, x2 + Inches(0.25), y + Inches(1.05),
             card_w - Inches(0.5), Inches(0.7),
             "Tangible Store", font=FONT_TITLE, size=28, bold=True, color=INK)
    add_text(s, x2 + Inches(0.25), y + Inches(1.75),
             card_w - Inches(0.5), Inches(0.4),
             "Boutique légale · achat définitif",
             font=FONT_BODY, size=14, italic=True, color=FLAME_DEEP)
    thin_line(s, x2 + Inches(0.25), y + Inches(2.2), Inches(0.7),
              color=FLAME, thickness_pt=1.5)
    store_items = [
        "Achat définitif (SD / HD / 4K)",
        "Certificat on-chain cryptographique",
        "Distribution P2P / IPFS",
        "Marché de revente avec royalties",
        "Paiement CB / crypto",
    ]
    for i, item in enumerate(store_items):
        by = y + Inches(2.45 + i * 0.3)
        add_text(s, x2 + Inches(0.3), by, card_w - Inches(0.5), Inches(0.3),
                 "✓  " + item, font=FONT_BODY, size=13, color=INK_SOFT)

    progress_footer(s, 3)
    set_notes(s, """
[1:00-2:30 — SOLUTION]

"Tangible, c'est la première plateforme où acheter un film signifie vraiment le posséder.
Un — Tangible Player. Media center local. Bibliothèque chiffrée. 100% hors-ligne.
Deux — Tangible Store. Boutique légale. Certificat de propriété cryptographique vérifiable hors-ligne. Revendable."
Ralentir sur "certificat" et "revendable" — mots nouveaux pour le jury.
""")

def slide_4_demo1(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY_DEEP)
    full_bg_gradient(s, IVORY, IVORY_DEEP, angle_deg=135)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "DÉMO · 1 / 2")
    big_title(s, Inches(0.95), "Achat en un clic.", size=44)
    underline_accent(s, Inches(0.7), Inches(2.1), length_in=1.2)

    ph_w = Inches(5.85)
    ph_h = Inches(4.5)
    y = Inches(2.4)

    # Left — Player
    left_card = rounded(s, Inches(0.7), y, ph_w, ph_h,
                        fill_color=PAPER, radius_in=0.2)
    set_shadow(left_card, blur_pt=30, dist_pt=10, alpha=25000)

    # Window chrome top
    tb1 = rounded(s, Inches(0.7), y, ph_w, Inches(0.45),
            fill_color=IVORY_DEEP, radius_in=0.2)
    for ci, cc in enumerate([RGBColor(0xE3,0x6A,0x5E),
                              RGBColor(0xE0,0xB0,0x4D),
                              RGBColor(0x5A,0xC0,0x5A)]):
        dot = add_rect(s, Inches(0.9) + Inches(0.3)*ci, y + Inches(0.14),
                 Inches(0.18), Inches(0.18),
                 fill_color=cc, shape=MSO_SHAPE.OVAL)
    add_text(s, Inches(2.0), y + Inches(0.08), Inches(4), Inches(0.3),
             "TANGIBLE PLAYER", font=FONT_MONO, size=10, bold=True, color=MUTED)
    add_text(s, Inches(0.7), y + Inches(1.8), ph_w, Inches(0.5),
             "[ Bibliothèque Player ]",
             font=FONT_MONO, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), y + Inches(2.25), ph_w, Inches(0.4),
             "remplacer par screenshot",
             font=FONT_BODY, size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), y + ph_h + Inches(0.25), ph_w, Inches(0.4),
             "📚  Bibliothèque chiffrée locale",
             font=FONT_BODY, size=14, bold=True, color=INK,
             align=PP_ALIGN.CENTER)

    # Right — Store
    x2 = Inches(0.7) + ph_w + Inches(0.2)
    right_card = rounded(s, x2, y, ph_w, ph_h,
                          fill_color=PAPER, radius_in=0.2)
    set_shadow(right_card, blur_pt=40, dist_pt=12,
               color=FLAME_DEEP, alpha=25000)
    # Store titlebar with gradient
    tb2 = rounded(s, x2, y, ph_w, Inches(0.45),
            fill_color=FLAME, radius_in=0.2)
    set_gradient_linear(tb2, FLAME_BRIGHT, FLAME_DEEP, angle_deg=0)
    add_text(s, x2 + Inches(0.2), y + Inches(0.08), Inches(4), Inches(0.3),
             "TANGIBLE STORE  ·  DUNE PART TWO",
             font=FONT_MONO, size=10, bold=True, color=IVORY)
    add_text(s, x2, y + Inches(1.8), ph_w, Inches(0.5),
             "[ Page film Store ]",
             font=FONT_MONO, size=13, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(s, x2, y + Inches(2.25), ph_w, Inches(0.4),
             "remplacer par screenshot",
             font=FONT_BODY, size=10, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

    # Price & CTA
    add_text(s, x2, y + Inches(3.1), ph_w, Inches(0.4),
             "14,99 €  ·  4K HDR",
             font=FONT_TITLE, size=20, bold=True, color=INK, align=PP_ALIGN.CENTER)
    # CTA with gradient + glow
    cta = rounded(s, x2 + Inches(1.3), y + Inches(3.55),
            ph_w - Inches(2.6), Inches(0.5),
            fill_color=FLAME, radius_in=0.08)
    set_gradient_linear(cta, FLAME_BRIGHT, FLAME_DEEP, angle_deg=90)
    set_shadow(cta, blur_pt=20, dist_pt=6, color=FLAME_DEEP, alpha=50000)
    add_text(s, x2 + Inches(1.3), y + Inches(3.58),
             ph_w - Inches(2.6), Inches(0.5),
             "ACHETER DÉFINITIVEMENT",
             font=FONT_MONO, size=11, bold=True, color=IVORY,
             align=PP_ALIGN.CENTER)

    add_text(s, x2, y + ph_h + Inches(0.25), ph_w, Inches(0.4),
             "🛒  Achat définitif + certificat",
             font=FONT_BODY, size=14, bold=True, color=FLAME_DEEP,
             align=PP_ALIGN.CENTER)

    progress_footer(s, 4)
    set_notes(s, """
[2:30-3:30 — DÉMO]
"Voici ce que ça donne. Deux écrans : à gauche, votre bibliothèque. À droite, la boutique."
⚠️ REMPLACER les placeholders par de vraies captures avant le jury.
""")

def slide_5_demo2(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INK)

    # Radial gradient dark
    full_bg_radial(s, INK_MID, INK_DEEP, focal_x=75000, focal_y=30000)

    icon = logos.get("tangible-icon")
    if icon:
        wm = watermark_logo(s, icon, size_in=8,
                            x_emu=Inches(10.5), y_emu=Inches(1.5),
                            alpha_amt=5000)
        if wm: set_soft_edge(wm, radius_pt=15)

    # Ambient glow behind certificate
    ambient_glow(s, Inches(4), Inches(4.5), radius_in=3.5,
                 color=FLAME, alpha=25000)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "DÉMO · 2 / 2")
    big_title(s, Inches(0.95), "La propriété, prouvée.",
              color=IVORY, size=44, italic=True)
    underline_accent(s, Inches(0.7), Inches(2.1), length_in=1.2)

    cx = Inches(0.7)
    cy = Inches(2.55)
    cw = Inches(6.4)
    ch = Inches(4.25)

    # Certificat with strong shadow + gradient
    cert = rounded(s, cx, cy, cw, ch, fill_color=PAPER, radius_in=0.18)
    set_gradient_linear(cert, PAPER, IVORY_DEEP, angle_deg=135)
    set_shadow(cert, blur_pt=50, dist_pt=14,
               color=RGBColor(0, 0, 0), alpha=55000)

    # Top ribbon with gradient
    ribbon = rounded(s, cx, cy, cw, Inches(0.13),
                     fill_color=FLAME, radius_in=0.04)
    set_gradient_linear(ribbon, FLAME_BRIGHT, FLAME_DEEP, angle_deg=0)
    set_glow(ribbon, radius_pt=12, color=FLAME, alpha=50000)

    add_text(s, cx + Inches(0.35), cy + Inches(0.35),
             cw - Inches(0.7), Inches(0.5),
             "🎟  CERTIFICAT DE PROPRIÉTÉ",
             font=FONT_MONO, size=11, bold=True, color=FLAME)
    add_text(s, cx + Inches(0.35), cy + Inches(0.75),
             cw - Inches(0.7), Inches(0.6),
             "Dune : Part Two",
             font=FONT_TITLE, size=26, bold=True, color=INK)
    add_text(s, cx + Inches(0.35), cy + Inches(1.25),
             cw - Inches(0.7), Inches(0.4),
             "Denis Villeneuve  ·  2024  ·  4K HDR Dolby Vision",
             font=FONT_BODY, size=12, italic=True, color=MUTED_DARK)

    thin_line(s, cx + Inches(0.35), cy + Inches(1.7),
              cw - Inches(0.7), thickness_pt=1, color=RGBColor(0xDD, 0xD8, 0xC5))

    details = [
        ("PROPRIÉTAIRE", "0xAB…E9"),
        ("LICENCE",      "TNGBL-D2-4K-0x9e3f…"),
        ("TX ON-CHAIN",  "0x7ea…d83"),
        ("ACHETÉ LE",    "15 avril 2026 · 14:32"),
    ]
    for i, (k, v) in enumerate(details):
        row = i // 2
        col = i % 2
        tx = cx + Inches(0.35 + col * 3.0)
        ty = cy + Inches(1.95 + row * 0.55)
        add_text(s, tx, ty, Inches(1.3), Inches(0.3),
                 k, font=FONT_MONO, size=8, bold=True, color=MUTED_DARK)
        add_text(s, tx, ty + Inches(0.25), Inches(3), Inches(0.3),
                 v, font=FONT_MONO, size=11, bold=True, color=INK)

    chip_y = cy + ch - Inches(0.7)
    chip1 = rounded(s, cx + Inches(0.35), chip_y, Inches(2.3), Inches(0.4),
            fill_color=RGBColor(0xE7, 0xF4, 0xEC), radius_in=0.1)
    add_text(s, cx + Inches(0.35), chip_y + Inches(0.04),
             Inches(2.3), Inches(0.4),
             "✓ VÉRIFIABLE HORS-LIGNE",
             font=FONT_MONO, size=10, bold=True, color=SUCCESS,
             align=PP_ALIGN.CENTER)
    chip2 = rounded(s, cx + Inches(2.8), chip_y, Inches(1.5), Inches(0.4),
            fill_color=RGBColor(0xFE, 0xEC, 0xDF), radius_in=0.1)
    add_text(s, cx + Inches(2.8), chip_y + Inches(0.04),
             Inches(1.5), Inches(0.4),
             "♻ REVENDABLE",
             font=FONT_MONO, size=10, bold=True, color=FLAME_DEEP,
             align=PP_ALIGN.CENTER)

    # Panel revente with gradient
    rx = Inches(7.45)
    ry = Inches(2.55)
    rw = Inches(5.15)
    rh = Inches(4.25)
    panel = rounded(s, rx, ry, rw, rh, fill_color=INK_DEEP, radius_in=0.18)
    set_gradient_linear(panel, INK, INK_DEEP, angle_deg=135)
    set_shadow(panel, blur_pt=40, dist_pt=10, alpha=50000)

    ribbon2 = rounded(s, rx, ry, rw, Inches(0.13), fill_color=FLAME, radius_in=0.04)
    set_gradient_linear(ribbon2, FLAME_BRIGHT, FLAME_DEEP, angle_deg=0)
    set_glow(ribbon2, radius_pt=12, color=FLAME, alpha=50000)

    add_text(s, rx + Inches(0.35), ry + Inches(0.35),
             rw - Inches(0.7), Inches(0.5),
             "♻  MARCHÉ SECONDAIRE",
             font=FONT_MONO, size=11, bold=True, color=FLAME)
    add_text(s, rx + Inches(0.35), ry + Inches(0.8),
             rw - Inches(0.7), Inches(0.9),
             "Revendre.\nLégalement.",
             font=FONT_TITLE, size=32, bold=True, italic=True, color=IVORY,
             line_spacing=1.0)
    add_text(s, rx + Inches(0.35), ry + Inches(2.0),
             rw - Inches(0.7), Inches(0.4),
             "EXEMPLE · VENTE À 8 €",
             font=FONT_MONO, size=10, bold=True, color=FLAME)
    thin_line(s, rx + Inches(0.35), ry + Inches(2.4),
              Inches(0.6), color=FLAME, thickness_pt=1)

    splits = [
        ("Vendeur",      80, "6,40 €"),
        ("Ayants droit", 15, "1,20 €"),
        ("Tangible",      5, "0,40 €"),
    ]
    bar_y = ry + Inches(2.7)
    for i, (who, pct, amount) in enumerate(splits):
        ty = bar_y + Inches(i * 0.48)
        add_text(s, rx + Inches(0.35), ty, Inches(2.2), Inches(0.3),
                 who, font=FONT_BODY, size=12, color=IVORY)
        max_bar = Inches(2.0)
        bw = Emu(int(max_bar * (pct / 100)))
        bg_bar = add_rect(s, rx + Inches(2.4), ty + Inches(0.07),
                          max_bar, Inches(0.2),
                          fill_color=RGBColor(0x2F, 0x34, 0x3A))
        bar = add_rect(s, rx + Inches(2.4), ty + Inches(0.07),
                        bw, Inches(0.2), fill_color=FLAME)
        set_gradient_linear(bar, FLAME_BRIGHT, FLAME_DEEP, angle_deg=0)
        set_glow(bar, radius_pt=6, color=FLAME, alpha=35000)
        add_text(s, rx + Inches(4.55), ty, Inches(0.7), Inches(0.3),
                 amount, font=FONT_MONO, size=11, bold=True, color=FLAME,
                 align=PP_ALIGN.RIGHT)

    progress_footer(s, 5, dark=True)
    set_notes(s, """
[3:30 — CERTIFICAT + REVENTE]
Pointer les 3 barres dans l'ordre. Le jury doit voir que le studio gagne AUSSI.
""")

def slide_6_difference(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    full_bg_gradient(s, PAPER, IVORY_DEEP, angle_deg=135)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "DIFFÉRENCIATION")
    big_title(s, Inches(0.95), "Tangible vs. le reste.", size=44)
    underline_accent(s, Inches(0.7), Inches(2.1), length_in=1.2)

    rows = [
        ("Achat définitif",         "✓", "~", "✕", "✕"),
        ("Téléchargement local",    "✓", "~", "~", "✓"),
        ("Certificat de propriété", "✓", "✕", "✕", "✕"),
        ("Revente possible",        "✓", "✕", "✕", "✕"),
        ("Hors-ligne à vie",        "✓", "✕", "✕", "✓"),
        ("Media center complet",    "✓", "~", "✕", "✓"),
        ("Catalogue légal",         "✓", "✓", "✓", "✕"),
        ("Pas de vendor lock",      "✓", "✕", "✕", "✓"),
    ]

    col0_w = Inches(3.8)
    col_w = Inches(2.05)
    table_x = Inches(0.8)
    table_y = Inches(2.55)
    header_h = Inches(0.7)
    row_h = Inches(0.42)

    xs = [table_x]
    xs.append(xs[-1] + col0_w)
    for _ in range(3):
        xs.append(xs[-1] + col_w)

    total_w = col0_w + col_w * 4
    total_h = header_h + row_h * len(rows)
    # Table bg with shadow
    tbl_bg = rounded(s, table_x, table_y, total_w, total_h,
                     fill_color=PAPER, radius_in=0.14)
    set_shadow(tbl_bg, blur_pt=40, dist_pt=10, alpha=25000)

    headers = ["Critère", "Tangible", "iTunes", "Netflix", "Jellyfin"]
    for i, h in enumerate(headers):
        w = col0_w if i == 0 else col_w
        fill = FLAME if i == 1 else (INK if i == 0 else INK_SOFT)
        if i == 0:
            hdr = rounded(s, xs[i], table_y, w, header_h, fill_color=fill, radius_in=0.14)
            set_gradient_linear(hdr, INK, INK_DEEP, angle_deg=180)
        elif i == 1:
            hdr = rounded(s, xs[i], table_y, w, header_h, fill_color=fill, radius_in=0.04)
            set_gradient_linear(hdr, FLAME_BRIGHT, FLAME_DEEP, angle_deg=180)
            set_glow(hdr, radius_pt=12, color=FLAME, alpha=40000)
        elif i == 4:
            hdr = rounded(s, xs[i], table_y, w, header_h, fill_color=fill, radius_in=0.14)
            set_gradient_linear(hdr, INK_SOFT, INK, angle_deg=180)
        else:
            hdr = add_rect(s, xs[i], table_y, w, header_h, fill_color=fill)
            set_gradient_linear(hdr, INK_SOFT, INK, angle_deg=180)
        add_text(s, xs[i], table_y + Inches(0.18), w, header_h,
                 h, font=FONT_BODY, size=15, bold=True, color=IVORY,
                 align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        y = table_y + header_h + row_h * ri
        bg = PAPER if ri % 2 == 0 else IVORY
        add_rect(s, xs[0], y, total_w, row_h, fill_color=bg)
        # Tangible col highlight
        add_rect(s, xs[1], y, col_w, row_h, fill_color=RGBColor(0xFD, 0xF4, 0xEC))
        add_text(s, xs[0] + Inches(0.3), y + Inches(0.08),
                 col0_w - Inches(0.3), row_h,
                 row[0], font=FONT_BODY, size=13, color=INK)
        for ci in range(1, 5):
            mark = row[ci]
            if mark == "✓":
                color, bold = SUCCESS, True
            elif mark == "~":
                color, bold = FLAME, True
            else:
                color, bold = DANGER, True
            add_text(s, xs[ci], y + Inches(0.06), col_w, row_h,
                     mark, font=FONT_MONO, size=18, bold=bold,
                     color=color, align=PP_ALIGN.CENTER)

    add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.4),
             "Une seule ligne avec tous les ✓ : la nôtre.",
             font=FONT_TITLE, size=17, bold=True, italic=True,
             color=FLAME, align=PP_ALIGN.CENTER)

    progress_footer(s, 6)
    set_notes(s, """
[3:30-4:30 — DIFFÉRENCIATION]
"iTunes ? Propriété fictive. Netflix ? Location. Jellyfin ? Pas de boutique légale. Tangible combine tout."
Laisser le jury lire 3 s en silence.
""")

def slide_7_marche(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    full_bg_radial(s, PAPER, IVORY_DEEP, focal_x=50000, focal_y=20000)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "LE MARCHÉ")
    big_title(s, Inches(0.95), "Un marché fatigué.", size=44)
    underline_accent(s, Inches(0.7), Inches(2.1), length_in=1.2)

    kpis = [
        ("130", "Mds €",  "Marché vidéo mondial",           "Statista 2024"),
        ("3,4", "",        "Abonnements SVOD / foyer FR",    "Médiamétrie 2024"),
        ("300", "€/an",    "Dépense moyenne par foyer",      "Deloitte 2024"),
    ]

    card_w = Inches(3.85)
    card_h = Inches(2.85)
    gap = Inches(0.2)
    y = Inches(2.55)
    for i, (big, unit, label, source) in enumerate(kpis):
        x = Inches(0.7) + (card_w + gap) * i
        card = rounded(s, x, y, card_w, card_h, fill_color=PAPER, radius_in=0.16)
        set_gradient_linear(card, PAPER, IVORY_DEEP, angle_deg=135)
        set_shadow(card, blur_pt=32, dist_pt=8, alpha=22000)
        # Top stripe gradient
        stripe = rounded(s, x, y, card_w, Inches(0.09),
                fill_color=FLAME, radius_in=0.03)
        set_gradient_linear(stripe, FLAME_BRIGHT, FLAME_DEEP, angle_deg=0)
        set_glow(stripe, radius_pt=10, color=FLAME, alpha=50000)
        # Big number
        add_text(s, x + Inches(0.3), y + Inches(0.45),
                 card_w - Inches(0.6), Inches(1.3),
                 big, font=FONT_TITLE, size=80, bold=True, color=INK,
                 line_spacing=1.0)
        if unit:
            add_text(s, x + Inches(0.3), y + Inches(1.55),
                     card_w - Inches(0.6), Inches(0.5),
                     unit, font=FONT_TITLE, size=22, bold=True, color=FLAME)
        add_text(s, x + Inches(0.3), y + Inches(2.05),
                 card_w - Inches(0.6), Inches(0.4),
                 label, font=FONT_BODY, size=13, color=INK_SOFT)
        add_text(s, x + Inches(0.3), y + card_h - Inches(0.4),
                 card_w - Inches(0.6), Inches(0.3),
                 source, font=FONT_MONO, size=9, color=MUTED)

    # Banner opportunity with gradient
    opp_y = Inches(5.65)
    banner = rounded(s, Inches(0.7), opp_y, Inches(11.9), Inches(1.3),
            fill_color=INK, radius_in=0.12)
    set_gradient_linear(banner, INK, INK_DEEP, angle_deg=90)
    set_shadow(banner, blur_pt=28, dist_pt=8, alpha=30000)
    # Left accent
    accent = rounded(s, Inches(0.7), opp_y, Inches(0.12), Inches(1.3),
            fill_color=FLAME, radius_in=0.03)
    set_gradient_linear(accent, FLAME_BRIGHT, FLAME_DEEP, angle_deg=90)
    set_glow(accent, radius_pt=10, color=FLAME, alpha=50000)

    add_text(s, Inches(1.1), opp_y + Inches(0.25),
             Inches(11), Inches(0.4),
             "NOTRE OPPORTUNITÉ  →",
             font=FONT_MONO, size=11, bold=True, color=FLAME)
    add_text(s, Inches(1.1), opp_y + Inches(0.6),
             Inches(11.5), Inches(0.7),
             "La niche propriété réelle  —  aucun acteur ne la sert correctement.",
             font=FONT_TITLE, size=22, bold=True, italic=True, color=IVORY)

    progress_footer(s, 7)
    set_notes(s, """
[4:30 — MARCHÉ]
Laisser un beat. Puis : "Personne ne sert la niche propriété réelle. C'est là qu'on entre."
""")

def slide_8_bm(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    full_bg_gradient(s, PAPER, IVORY_DEEP, angle_deg=135)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "BUSINESS MODEL")
    big_title(s, Inches(0.95), "Quatre revenus. Diversifiés.", size=40)
    underline_accent(s, Inches(0.7), Inches(2.1), length_in=1.2)

    revenues = [
        ("20 %",      "Commission vente",    "Studio 70 · Seeders 10 · Tangible 20", True),
        ("5 %",       "Commission revente",  "Vendeur 80 · Ayants droit 15 · Tangible 5", False),
        ("7,99 €",    "Tangible Pass / mois","Catalogue rotatif · −15 % sur achats", False),
        ("B2B",       "Licence SDK",         "Studios · distributeurs · API tierce", False),
    ]

    card_w = Inches(5.95)
    card_h = Inches(2.0)
    y_start = Inches(2.5)

    for i, (big, title, detail, primary) in enumerate(revenues):
        row = i // 2
        col = i % 2
        x = Inches(0.7) + (card_w + Inches(0.2)) * col
        y = y_start + (card_h + Inches(0.25)) * row
        card = rounded(s, x, y, card_w, card_h,
                       fill_color=PAPER, radius_in=0.14)
        set_gradient_linear(card, PAPER, IVORY_DEEP, angle_deg=135)
        if primary:
            set_shadow(card, blur_pt=36, dist_pt=10, color=FLAME_DEEP, alpha=30000)
        else:
            set_shadow(card, blur_pt=26, dist_pt=7, alpha=18000)
        # Left stripe
        stripe = rounded(s, x, y, Inches(0.14), card_h,
                fill_color=FLAME if primary else INK, radius_in=0.05)
        if primary:
            set_gradient_linear(stripe, FLAME_BRIGHT, FLAME_DEEP, angle_deg=90)
            set_glow(stripe, radius_pt=12, color=FLAME, alpha=55000)
        # Big
        add_text(s, x + Inches(0.4), y + Inches(0.3),
                 Inches(2.3), Inches(1.5),
                 big, font=FONT_TITLE, size=34, bold=True, color=FLAME)
        add_text(s, x + Inches(2.8), y + Inches(0.35),
                 card_w - Inches(3), Inches(0.55),
                 title, font=FONT_TITLE, size=19, bold=True, color=INK)
        add_text(s, x + Inches(2.8), y + Inches(0.95),
                 card_w - Inches(3), Inches(0.9),
                 detail, font=FONT_BODY, size=12, color=INK_SOFT,
                 line_spacing=1.35)
        if primary:
            badge = rounded(s, x + card_w - Inches(1.2), y + Inches(0.3),
                    Inches(0.9), Inches(0.35),
                    fill_color=FLAME, radius_in=0.1)
            set_gradient_linear(badge, FLAME_BRIGHT, FLAME_DEEP, angle_deg=90)
            set_glow(badge, radius_pt=8, color=FLAME, alpha=50000)
            add_text(s, x + card_w - Inches(1.2), y + Inches(0.33),
                     Inches(0.9), Inches(0.35),
                     "PRIMARY", font=FONT_MONO, size=9, bold=True, color=IVORY,
                     align=PP_ALIGN.CENTER)

    progress_footer(s, 8)
    set_notes(s, """
[4:30-5:30 — BUSINESS MODEL]
Quatre revenus diversifiés. Insister que le studio gagne AUSSI sur la revente (15%).
""")

def slide_9_financier(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    full_bg_gradient(s, PAPER, IVORY_DEEP, angle_deg=135)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "FINANCIER")
    big_title(s, Inches(0.95), "Break-even en Année 3.", size=40)
    underline_accent(s, Inches(0.7), Inches(2.1), length_in=1.2)

    chart_x = Inches(0.7)
    chart_y = Inches(2.55)
    chart_w = Inches(7.8)
    chart_h = Inches(4.1)

    chart_card = rounded(s, chart_x, chart_y, chart_w, chart_h,
                          fill_color=PAPER, radius_in=0.16)
    set_gradient_linear(chart_card, PAPER, IVORY_DEEP, angle_deg=135)
    set_shadow(chart_card, blur_pt=32, dist_pt=8, alpha=22000)

    add_text(s, chart_x + Inches(0.4), chart_y + Inches(0.25),
             chart_w - Inches(0.8), Inches(0.4),
             "CHIFFRE D'AFFAIRES (k€)",
             font=FONT_MONO, size=10, bold=True, color=MUTED_DARK)
    add_text(s, chart_x + Inches(0.4), chart_y + Inches(0.6),
             chart_w - Inches(0.8), Inches(0.4),
             "×33 en 3 ans",
             font=FONT_TITLE, size=22, bold=True, italic=True, color=FLAME)

    plot_x = chart_x + Inches(0.8)
    plot_y = chart_y + Inches(1.3)
    plot_w = chart_w - Inches(1.2)
    plot_h = chart_h - Inches(2.2)

    grid_color = RGBColor(0xE8, 0xE2, 0xD3)
    max_val = 7000
    for g in range(5):
        gy = plot_y + Emu(int(plot_h * (g / 4)))
        add_rect(s, plot_x, gy, plot_w, Emu(int(Pt(0.5))),
                 fill_color=grid_color)
        val = int(max_val * (1 - g / 4))
        add_text(s, plot_x - Inches(0.75), gy - Inches(0.11),
                 Inches(0.7), Inches(0.3),
                 f"{val//1000 if val >= 1000 else val}{'k' if val < 1000 else 'M'}",
                 font=FONT_MONO, size=9, color=MUTED, align=PP_ALIGN.RIGHT)

    bars = [
        ("2026", 200,  "200 k€", False),
        ("2027", 1600, "1,6 M€", False),
        ("2028", 6700, "6,7 M€", True),
    ]
    bar_w = Inches(1.4)
    slot = plot_w / len(bars)
    for i, (year, val, label, hero) in enumerate(bars):
        bh = Emu(int(plot_h * (val / max_val)))
        bx = plot_x + slot * i + (slot - bar_w) // 2
        by = plot_y + plot_h - bh
        # gradient bar
        bar = rounded(s, bx, by, bar_w, bh,
                      fill_color=FLAME, radius_in=0.06)
        if hero:
            set_gradient_linear(bar, FLAME_BRIGHT, FLAME_DEEP, angle_deg=90)
            set_glow(bar, radius_pt=16, color=FLAME, alpha=45000)
            set_shadow(bar, blur_pt=20, dist_pt=4, color=FLAME_DEEP, alpha=40000)
        elif i == 1:
            set_gradient_linear(bar, RGBColor(0xF5, 0xA5, 0x78),
                                RGBColor(0xC5, 0x58, 0x2A), angle_deg=90)
            set_shadow(bar, blur_pt=14, dist_pt=4, alpha=25000)
        else:
            set_gradient_linear(bar, RGBColor(0xF8, 0xCD, 0xAE),
                                RGBColor(0xD8, 0x9E, 0x74), angle_deg=90)
            set_shadow(bar, blur_pt=10, dist_pt=3, alpha=20000)

        add_text(s, bx - Inches(0.2), by - Inches(0.55),
                 bar_w + Inches(0.4), Inches(0.5),
                 label, font=FONT_TITLE, size=17, bold=True, color=INK,
                 align=PP_ALIGN.CENTER)
        add_text(s, bx - Inches(0.2), plot_y + plot_h + Inches(0.15),
                 bar_w + Inches(0.4), Inches(0.4),
                 year, font=FONT_MONO, size=13, bold=True, color=INK_SOFT,
                 align=PP_ALIGN.CENTER)

    be_x = plot_x + slot * 2 + slot / 2
    be_y = plot_y - Inches(0.05)
    add_text(s, be_x - Inches(1.4), be_y - Inches(0.35),
             Inches(2.8), Inches(0.3),
             "▼ BREAK-EVEN",
             font=FONT_MONO, size=9, bold=True, color=SUCCESS,
             align=PP_ALIGN.CENTER)

    # Ask panel
    ax = Inches(8.7)
    ay = Inches(2.55)
    aw = Inches(3.95)
    ah = Inches(4.1)
    ask_panel = rounded(s, ax, ay, aw, ah,
                        fill_color=INK_DEEP, radius_in=0.18)
    set_gradient_linear(ask_panel, INK, INK_DEEP, angle_deg=135)
    set_shadow(ask_panel, blur_pt=45, dist_pt=12, alpha=50000)

    # Ambient glow inside ask panel behind 950
    ambient_glow(s, ax + aw/2, ay + Inches(1.3), radius_in=1.6,
                 color=FLAME, alpha=28000)

    # Top stripe gradient
    stripe = rounded(s, ax, ay, aw, Inches(0.12), fill_color=FLAME, radius_in=0.04)
    set_gradient_linear(stripe, FLAME_BRIGHT, FLAME_DEEP, angle_deg=0)
    set_glow(stripe, radius_pt=14, color=FLAME, alpha=55000)

    add_text(s, ax + Inches(0.35), ay + Inches(0.3),
             aw - Inches(0.7), Inches(0.4),
             "AMORÇAGE DEMANDÉ",
             font=FONT_MONO, size=11, bold=True, color=FLAME)
    add_text(s, ax + Inches(0.35), ay + Inches(0.85),
             aw - Inches(0.7), Inches(1.4),
             "950 k€",
             font=FONT_TITLE, size=72, bold=True, color=IVORY,
             line_spacing=1.0)
    add_text(s, ax + Inches(0.35), ay + Inches(2.35),
             aw - Inches(0.7), Inches(0.5),
             "500 k€ seed  +  450 k€ prêt & BPI",
             font=FONT_BODY, size=13, color=MUTED_LIGHT)

    thin_line(s, ax + Inches(0.35), ay + Inches(2.95),
              Inches(0.7), color=FLAME, thickness_pt=1.5)

    usage = ["12 mois de runway", "MVP Player", "Store bêta", "5 000 users"]
    for i, u in enumerate(usage):
        add_text(s, ax + Inches(0.35), ay + Inches(3.1 + i * 0.24),
                 aw - Inches(0.7), Inches(0.3),
                 "→  " + u, font=FONT_BODY, size=11, color=IVORY)

    progress_footer(s, 9)
    set_notes(s, """
[5:00 — FINANCIER + ASK]
Pointer la barre orange 2028. Puis les yeux vers le panel 950 k€.
""")

def slide_10_equipe_roadmap(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, IVORY)
    full_bg_radial(s, PAPER, IVORY_DEEP, focal_x=50000, focal_y=20000)

    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.5), "ÉQUIPE & ROADMAP")
    big_title(s, Inches(0.95),
              "Des fondateurs. Une trajectoire.", size=36)
    underline_accent(s, Inches(0.7), Inches(2.05), length_in=1.2)

    team_x = Inches(0.7)
    team_y = Inches(2.45)
    add_text(s, team_x, team_y, Inches(5.8), Inches(0.5),
             "L'ÉQUIPE",
             font=FONT_MONO, size=11, bold=True, color=FLAME)
    thin_line(s, team_x, team_y + Inches(0.45), Inches(0.7),
              color=FLAME, thickness_pt=1.5)

    team = [
        ("CEO",  "Produit & Business",      "Vision · levée · studios", True),
        ("CTO",  "Architecture & Dev",       "Stack · qualité · recrutement", False),
        ("CSO",  "Sécurité & Cryptographie", "5 couches · audits · conformité", False),
        ("CPO",  "UX & Design",              "Player · Store · brand", False),
    ]
    for i, (role, title, detail, highlight) in enumerate(team):
        y = team_y + Inches(0.85 + i * 0.85)
        badge = rounded(s, team_x, y, Inches(0.9), Inches(0.6),
                        fill_color=FLAME if highlight else INK, radius_in=0.1)
        if highlight:
            set_gradient_linear(badge, FLAME_BRIGHT, FLAME_DEEP, angle_deg=135)
            set_glow(badge, radius_pt=10, color=FLAME, alpha=50000)
        else:
            set_gradient_linear(badge, INK_SOFT, INK, angle_deg=135)
        add_text(s, team_x, y + Inches(0.14), Inches(0.9), Inches(0.4),
                 role, font=FONT_MONO, size=11, bold=True, color=IVORY,
                 align=PP_ALIGN.CENTER)
        add_text(s, team_x + Inches(1.1), y + Inches(0.05),
                 Inches(4.5), Inches(0.4),
                 title, font=FONT_TITLE, size=16, bold=True, color=INK)
        add_text(s, team_x + Inches(1.1), y + Inches(0.38),
                 Inches(4.5), Inches(0.35),
                 detail, font=FONT_BODY, size=11, italic=True, color=MUTED_DARK)

    road_x = Inches(7)
    road_y = Inches(2.45)
    add_text(s, road_x, road_y, Inches(5.8), Inches(0.5),
             "LA ROADMAP",
             font=FONT_MONO, size=11, bold=True, color=FLAME)
    thin_line(s, road_x, road_y + Inches(0.45), Inches(0.7),
              color=FLAME, thickness_pt=1.5)

    thin_line(s, road_x + Inches(0.4), road_y + Inches(0.9),
              Inches(3.8), vertical=True, color=INK_LINE, thickness_pt=1.5)

    phases = [
        ("Phase 1",  "0-6 mo",    "MVP Player · 50 films indé", True),
        ("Phase 2",  "6-12 mo",   "Store · mobile · blockchain", False),
        ("Phase 3",  "12-24 mo",  "Marché secondaire · majors · B2B", False),
        ("Phase 4",  "24+ mo",    "Smart TV · international", False),
    ]
    for i, (phase, dur, detail, active) in enumerate(phases):
        y = road_y + Inches(1.0 + i * 0.95)
        dot = add_rect(s, road_x + Inches(0.32), y + Inches(0.2),
                 Inches(0.22), Inches(0.22),
                 fill_color=FLAME if active else INK,
                 shape=MSO_SHAPE.OVAL)
        if active:
            set_gradient_radial(dot, FLAME_BRIGHT, FLAME_DEEP)
            set_glow(dot, radius_pt=10, color=FLAME, alpha=60000)
        add_text(s, road_x + Inches(0.85), y,
                 Inches(5), Inches(0.4),
                 phase, font=FONT_MONO, size=11, bold=True, color=FLAME)
        add_text(s, road_x + Inches(2.0), y,
                 Inches(4), Inches(0.4),
                 "·  " + dur, font=FONT_MONO, size=10, color=MUTED)
        add_text(s, road_x + Inches(0.85), y + Inches(0.35),
                 Inches(5), Inches(0.4),
                 detail, font=FONT_TITLE, size=15, bold=True, color=INK)

    progress_footer(s, 10)
    set_notes(s, """
[5:30-6:15 — ÉQUIPE & ROADMAP]
Pointer la dot orange Phase 1 : "c'est ici qu'on commence."
""")

def slide_11_ask(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INK_DEEP)
    # Radial ambient
    full_bg_radial(s, INK, INK_DEEP, focal_x=25000, focal_y=40000)

    # Big ambient glow left
    ambient_glow(s, Inches(2.8), Inches(3.5), radius_in=5,
                 color=FLAME, alpha=32000)

    # Giant "950" decoration
    add_text(s, Inches(4.5), Inches(0.6), Inches(9), Inches(7),
             "950", font=FONT_TITLE, size=600, bold=True,
             color=RGBColor(0x1B, 0x1D, 0x21),
             align=PP_ALIGN.CENTER)

    icon = logos.get("tangible-icon")
    if icon:
        add_logo_picture(s, icon, Inches(12.2), Inches(0.45), 0.9)

    corner_mark(s, Inches(12.0), Inches(0.4), color=FLAME)
    eyebrow(s, Inches(0.7), Inches(0.55), "NOTRE ASK", color=FLAME)

    big_title(s, Inches(1.0),
              "950 k€ pour changer\nla propriété numérique.",
              color=IVORY, size=44, italic=False)
    underline_accent(s, Inches(0.7), Inches(3.1), length_in=1.5)

    cols = [
        ("500 k€",  "Seed",        "Dev, équipe\nfondatrice, MVP", True),
        ("450 k€",  "Prêt + BPI",  "Fonds de roulement,\nR&D sécurité", False),
        ("12 mois", "de runway",   "Jusqu'à 5 000 users\nvalidés", False),
    ]
    card_w = Inches(3.95)
    card_h = Inches(2.6)
    gap = Inches(0.2)
    y = Inches(3.6)
    for i, (big, label, detail, hero) in enumerate(cols):
        x = Inches(0.7) + (card_w + gap) * i
        card = rounded(s, x, y, card_w, card_h,
                       fill_color=INK_MID, radius_in=0.16)
        set_gradient_linear(card, INK_MID, INK_DEEP, angle_deg=135)
        if hero:
            set_shadow(card, blur_pt=40, dist_pt=10, color=FLAME_DEEP, alpha=40000)
        else:
            set_shadow(card, blur_pt=30, dist_pt=8, alpha=50000)
        # Top stripe
        stripe = rounded(s, x, y, card_w, Inches(0.09),
                fill_color=FLAME if hero else INK_LINE, radius_in=0.03)
        if hero:
            set_gradient_linear(stripe, FLAME_BRIGHT, FLAME_DEEP, angle_deg=0)
            set_glow(stripe, radius_pt=14, color=FLAME, alpha=55000)

        add_text(s, x + Inches(0.35), y + Inches(0.35),
                 card_w - Inches(0.7), Inches(0.9),
                 big, font=FONT_TITLE, size=36, bold=True, color=FLAME,
                 line_spacing=1.0)
        add_text(s, x + Inches(0.35), y + Inches(1.3),
                 card_w - Inches(0.7), Inches(0.4),
                 label, font=FONT_BODY, size=14, bold=True, color=IVORY)
        add_text(s, x + Inches(0.35), y + Inches(1.75),
                 card_w - Inches(0.7), Inches(0.8),
                 detail, font=FONT_BODY, size=12, color=MUTED_LIGHT,
                 line_spacing=1.3)

    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.5),
             "Au-delà de l'argent : une conviction.",
             font=FONT_TITLE, size=19, italic=True, color=MUTED_LIGHT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.75), Inches(12), Inches(0.4),
             "le numérique responsable passe par la propriété réelle.",
             font=FONT_BODY, size=13, italic=True,
             color=RGBColor(0xB5, 0xB0, 0xA5), align=PP_ALIGN.CENTER)

    progress_footer(s, 11, dark=True)
    set_notes(s, """
[6:15-7:00 — ASK + VISION]
SILENCE d'1 seconde avant "Mais au-delà de l'argent".
""")

def slide_12_closing(prs, logos):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s, INK_DEEP)

    # Gradient full bg
    full_bg_gradient(s, INK, INK_DEEP, angle_deg=180)
    # Ambient glow bottom center (behind slogan)
    ambient_glow(s, Inches(13.333/2), Inches(5.2), radius_in=4.5,
                 color=FLAME, alpha=28000)
    # Big halo behind logo
    ambient_glow(s, Inches(13.333/2), Inches(2.6), radius_in=3,
                 color=FLAME, alpha=35000)

    # Corner marks both
    corner_mark(s, Inches(13.0), Inches(0.35), color=FLAME)
    t = Emu(int(Pt(1.5)))
    length = Inches(0.3)
    add_rect(s, Inches(0.35), Inches(0.35), length, t, fill_color=FLAME)
    add_rect(s, Inches(0.35), Inches(0.35), t, length, fill_color=FLAME)

    logo = logos.get("tangible-logo-vertical")
    if logo:
        pic, lw, _ = add_logo_picture(s, logo, Inches(0), Inches(1.0), 3.2)
        s.shapes[-1].left = Inches(13.333/2) - lw // 2

    thin_line(s, Inches(13.333/2 - 1.0), Inches(4.55),
              Inches(2.0), thickness_pt=1.5, color=FLAME)

    add_text(s, Inches(0.7), Inches(4.85), Inches(12), Inches(1.0),
             "Ne louez plus votre passion. Possédez-la.",
             font=FONT_TITLE, size=44, bold=True, italic=True,
             color=IVORY, align=PP_ALIGN.CENTER, line_spacing=1.05)
    add_text(s, Inches(0.7), Inches(5.95), Inches(12), Inches(0.4),
             "·  DON'T RENT YOUR PASSION. OWN IT.  ·",
             font=FONT_MONO, size=13, bold=True, color=FLAME,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(6.6), Inches(12), Inches(0.4),
             "pixnop.github.io/tangible",
             font=FONT_MONO, size=12, color=MUTED_LIGHT,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.7), Inches(7.0), Inches(12), Inches(0.4),
             "— Merci. Vos questions. —",
             font=FONT_BODY, size=13, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)

    set_notes(s, """
[FIN — 7:00]
SILENCE. Regarder le jury. Attendre les questions.
Q&R prêtes : content/05 - Pitch/Objections et Réponses.md (20 Q/R).
""")

# ─────────────────── MAIN ──────────────────
def main():
    print("→ Préparation des logos...")
    logos = prepare_logos()
    print(f"  {len(logos)} logos OK")

    print("→ Génération (avec gradients, glows, shadows)...")
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_1_cover,
        slide_2_probleme,
        slide_3_solution,
        slide_4_demo1,
        slide_5_demo2,
        slide_6_difference,
        slide_7_marche,
        slide_8_bm,
        slide_9_financier,
        slide_10_equipe_roadmap,
        slide_11_ask,
        slide_12_closing,
    ]
    for i, build in enumerate(builders, 1):
        build(prs, logos)
        print(f"  Slide {i:2d} ✓  {build.__name__}")

    prs.save(OUTPUT)
    size_kb = OUTPUT.stat().st_size / 1024
    print(f"\n✅ {OUTPUT.relative_to(ROOT)}  ({size_kb:.0f} KB · 12 slides)")
    print("   → Gradients · Glows · Drop shadows · Soft edges appliqués")

if __name__ == "__main__":
    main()
